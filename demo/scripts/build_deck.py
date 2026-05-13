"""Build demo/CollectMind_Interview_Deck.pptx via python-pptx.

15 slides. Light background, violet-500 accent, Inter/Segoe sans + JetBrains
Mono/Consolas. Every claim slide carries a small-text footer citing the
backing artifact (ADR id, readiness-review section, principle number).

Run: python demo/scripts/build_deck.py
Output: demo/CollectMind_Interview_Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
OUT = REPO_ROOT / "demo" / "CollectMind_Interview_Deck.pptx"

# --- design tokens --------------------------------------------------------
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)        # violet-500
ACCENT_DARK = RGBColor(0x6D, 0x28, 0xD9)   # violet-700
INK = RGBColor(0x18, 0x18, 0x1B)           # zinc-950
INK_SOFT = RGBColor(0x3F, 0x3F, 0x46)      # zinc-700
INK_MUTED = RGBColor(0x71, 0x71, 0x7A)     # zinc-500
INK_FAINT = RGBColor(0xA1, 0xA1, 0xAA)     # zinc-400
BG = RGBColor(0xFA, 0xFA, 0xFA)            # zinc-50
PANEL = RGBColor(0xF4, 0xF4, 0xF5)         # zinc-100
PANEL_BORDER = RGBColor(0xE4, 0xE4, 0xE7)  # zinc-200
OK = RGBColor(0x05, 0x96, 0x69)            # emerald-600
WARN = RGBColor(0xD9, 0x77, 0x06)          # amber-600

SANS = "Inter"
SANS_FALLBACK = "Segoe UI"
MONO = "JetBrains Mono"
MONO_FALLBACK = "Consolas"


# --- helpers --------------------------------------------------------------
def add_paragraph(
    text_frame,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    font: str = SANS,
    mono: bool = False,
    align: int | None = None,
    space_before: int = 0,
    space_after: int = 0,
):
    p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    run = p.runs[0]
    run.font.name = MONO if mono else font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    mono: bool = False,
    align: int | None = None,
    word_wrap: bool = True,
):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    if text:
        add_paragraph(tf, text, size=size, bold=bold, color=color, mono=mono, align=align)
    return tx


def add_filled_rect(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: RGBColor | None = None,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_panel(slide, *, left, top, width, height, fill=PANEL, border=PANEL_BORDER):
    return add_filled_rect(
        slide, left=left, top=top, width=width, height=height, fill=fill, line=border
    )


def add_accent_bar(slide, *, top: float = 0.0):
    add_filled_rect(slide, left=0, top=top, width=13.333, height=0.06, fill=ACCENT)


def add_slide_title(slide, kicker: str, title: str, footer: str | None = None):
    add_accent_bar(slide)
    tx_kicker = add_textbox(
        slide,
        left=0.6,
        top=0.35,
        width=12.0,
        height=0.3,
        text=kicker.upper(),
        size=10,
        bold=True,
        color=ACCENT_DARK,
    )
    # spacing
    add_textbox(
        slide,
        left=0.6,
        top=0.68,
        width=12.0,
        height=0.7,
        text=title,
        size=30,
        bold=True,
        color=INK,
    )
    if footer:
        add_textbox(
            slide,
            left=0.6,
            top=6.85,
            width=12.0,
            height=0.35,
            text=footer,
            size=10,
            color=INK_MUTED,
            mono=True,
        )
    return tx_kicker


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


# --- slide builders -------------------------------------------------------
def slide_title(prs):
    slide = new_slide(prs)
    add_accent_bar(slide, top=0.0)
    # Mark / logo block
    add_filled_rect(
        slide, left=0.6, top=1.2, width=0.55, height=0.55, fill=ACCENT, line=None
    )
    add_textbox(
        slide,
        left=0.6,
        top=1.2,
        width=0.55,
        height=0.55,
        text="◆",
        size=28,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        left=1.3,
        top=1.18,
        width=10.0,
        height=0.7,
        text="CollectMind",
        size=44,
        bold=True,
        color=INK,
    )
    add_textbox(
        slide,
        left=1.3,
        top=1.85,
        width=11.0,
        height=0.45,
        text="Agentic vehicle-telemetry policy engine",
        size=18,
        color=INK_SOFT,
    )

    # Hero one-liner
    add_textbox(
        slide,
        left=0.6,
        top=3.1,
        width=12.0,
        height=1.2,
        text="The diagnostic-to-collection loop, closed under SLM rigor.",
        size=32,
        bold=True,
        color=INK,
    )
    add_textbox(
        slide,
        left=0.6,
        top=4.4,
        width=12.0,
        height=0.5,
        text="Open-weight SLM under decode-time grammar constraints · "
        "multi-tenant RLS · immutable audit chain · deterministic CI budget.",
        size=16,
        color=INK_SOFT,
    )

    # Author + date
    add_textbox(
        slide,
        left=0.6,
        top=6.4,
        width=12.0,
        height=0.4,
        text="Arun Veligatla · Sonatus interview demo · 2026-05-14",
        size=14,
        color=INK_MUTED,
    )
    add_textbox(
        slide,
        left=0.6,
        top=6.85,
        width=12.0,
        height=0.3,
        text="arunveligatla99.github.io/snts_collectmind  ·  github.com/arunveligatla99/snts_collectmind",
        size=10,
        color=INK_MUTED,
        mono=True,
    )
    return slide


def slide_gap(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "The gap",
        "Between an AI Technician hypothesis and a deployed collection policy",
        footer="constitution v1.0.1 · Principle I (production-grade by default)",
    )

    # Three-row "what gets lost" panel
    rows = [
        (
            "Today",
            "Diagnostic hypothesis → Slack thread → ticket → engineer hand-codes a policy → reviewer eyeballs it → ops deploys.",
        ),
        (
            "Lost in the handoff",
            "Provenance (who generated, which model, which seed, which template). "
            "Multi-tenant scoping. Outcome feedback. The fingerprint of the deployed policy.",
        ),
        (
            "Cost",
            "Audit is reconstructed from logs, not queried. Tenant isolation is policy, not enforcement. "
            "Outcome data never makes it back to the next hypothesis.",
        ),
    ]
    y = 1.7
    for label, body in rows:
        add_panel(slide, left=0.6, top=y, width=12.0, height=1.5)
        add_textbox(
            slide,
            left=0.85,
            top=y + 0.18,
            width=11.5,
            height=0.4,
            text=label.upper(),
            size=11,
            bold=True,
            color=ACCENT_DARK,
        )
        add_textbox(
            slide,
            left=0.85,
            top=y + 0.55,
            width=11.5,
            height=0.95,
            text=body,
            size=16,
            color=INK,
        )
        y += 1.7

    return slide


def slide_what_it_does(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "What CollectMind does",
        "Closes the loop with four agentic nodes + an immutable audit backbone",
        footer="constitution Principle XII (agent boundaries) · ADR-0003 (constrained decoding)",
    )

    nodes = [
        ("Orchestrator", "Deterministic Python. Routes the finding through the graph."),
        ("Policy Generator", "Only node that calls a model. SLM emits a CollectionPolicySpec under outlines grammar constraints."),
        ("Policy Validator", "VSS v6.0 vocabulary + PII consent + governance flags. Failures route back to the generator with errors in the retry prompt."),
        ("Policy Deployer", "Signs at registry-write. Ships through the Collector AI client. Tenant-scope check is the FIRST gate."),
        ("Audit backbone", "Every node writes an immutable audit row. FR-017a minimum field set on every kind=generated row."),
    ]
    y = 1.7
    for label, body in nodes:
        add_panel(slide, left=0.6, top=y, width=12.0, height=0.9)
        add_textbox(
            slide,
            left=0.85,
            top=y + 0.16,
            width=4.0,
            height=0.55,
            text=label,
            size=16,
            bold=True,
            color=ACCENT_DARK,
        )
        add_textbox(
            slide,
            left=4.9,
            top=y + 0.16,
            width=7.6,
            height=0.6,
            text=body,
            size=14,
            color=INK,
        )
        y += 1.05

    return slide


def slide_architecture(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Architecture",
        "Cloud control plane · SLM container has no egress except OTLP",
        footer="constitution Principle XIII (SLM-first, isolated, swappable) · `infra/terraform/networking/main.tf`",
    )

    # Boxes laid out as a left-to-right flow with audit_events as a band
    def box(label, sub, left, top, w=2.05, h=1.0, fill=PANEL):
        add_panel(slide, left=left, top=top, width=w, height=h, fill=fill)
        add_textbox(
            slide,
            left=left + 0.1,
            top=top + 0.1,
            width=w - 0.2,
            height=0.4,
            text=label,
            size=12,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left=left + 0.1,
            top=top + 0.5,
            width=w - 0.2,
            height=h - 0.5,
            text=sub,
            size=9,
            color=INK_SOFT,
            align=PP_ALIGN.CENTER,
        )

    def arrow(x1, y1, x2, y2, color=INK_FAINT):
        line = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = color
        line.line.width = Pt(1.5)
        return line

    top1 = 2.0
    top2 = 4.2
    # row 1
    box("AI Technician", "diagnostic findings", 0.4, top1)
    box("Orchestration API", "POST /findings · JWT", 2.65, top1)
    box("Rate limiter", "token bucket · Redis", 4.9, top1)
    box("LangGraph", "Orch · Gen · Validate · Deploy", 7.15, top1)
    box("SLM container", "Qwen2.5-7B · vLLM · outlines", 9.4, top1, fill=RGBColor(0xEE, 0xE9, 0xFE))
    # row 2
    box("Query API", "/policies /audit /outcome", 0.4, top2)
    box("Audit-Admin", "break-glass · operator JWT", 2.65, top2)
    box("Registry", "Postgres + Timescale · RLS", 4.9, top2)
    box("Collector AI", "deployment client", 7.15, top2)
    box("Observability", "OTel · Prom · Grafana", 9.4, top2, fill=RGBColor(0xEE, 0xE9, 0xFE))

    # arrows row1
    arrow(2.45, top1 + 0.5, 2.65, top1 + 0.5)
    arrow(4.70, top1 + 0.5, 4.90, top1 + 0.5)
    arrow(6.95, top1 + 0.5, 7.15, top1 + 0.5)
    arrow(9.20, top1 + 0.5, 9.40, top1 + 0.5)
    # row1 -> row2 (LangGraph → Registry; LangGraph → Collector; SLM → Observability)
    arrow(8.18, top1 + 1.0, 5.92, top2)
    arrow(8.18, top1 + 1.0, 8.18, top2)
    arrow(10.42, top1 + 1.0, 10.42, top2)

    # audit band
    add_filled_rect(
        slide, left=0.4, top=5.7, width=11.0, height=0.55, fill=ACCENT, line=None
    )
    add_textbox(
        slide,
        left=0.4,
        top=5.72,
        width=11.0,
        height=0.5,
        text="audit_events  —  every node writes an immutable row · FR-017a fields on every kind=generated",
        size=12,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    return slide


def slide_constitution(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "The constitution",
        "18 principles · 7 non-negotiable · v1.0.1 (ratified 2026-05-09)",
        footer=".specify/memory/constitution.md · NON-NEGOTIABLE deviations require a constitution amendment, not an ADR",
    )

    nn = [
        ("IV", "Tests Are Load-Bearing", "≥85% coverage · test-first · red-phase commits"),
        ("VII", "CI/CD Gates Merges", "lint · type · unit · contract · integration · scan · SBOM"),
        ("IX", "Security as a First-Class Requirement", "JWT · Pydantic v2 · SHA-pinned weights · gitleaks"),
        ("X", "Vehicle Telemetry Data Handling", "VSS v6.0 · PII consent · per-tenant isolation · 90-day retention"),
        ("XI", "Performance SLOs Are Measured, Not Aspired", "binding SC-001…SC-014 · breach fails the build"),
        ("XIII", "SLM-First, Isolated, Swappable", "open-weight SLM default · LLM is opt-in · pinned revision SHA"),
        ("XIV", "Deterministic, Budgeted Model Execution in CI", "real SLM in contract+int · stub in load+soak · gated full"),
    ]

    cols = 2
    cell_w = 5.95
    cell_h = 1.05
    start_x = 0.6
    start_y = 1.9
    gap_x = 0.2
    gap_y = 0.15
    for i, (num, label, body) in enumerate(nn):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        add_panel(slide, left=x, top=y, width=cell_w, height=cell_h)
        add_textbox(
            slide, left=x + 0.18, top=y + 0.12, width=0.6, height=0.45,
            text=num, size=20, bold=True, color=ACCENT_DARK, mono=True,
        )
        add_textbox(
            slide, left=x + 0.85, top=y + 0.10, width=cell_w - 1.0, height=0.4,
            text=label, size=14, bold=True, color=INK,
        )
        add_textbox(
            slide, left=x + 0.85, top=y + 0.48, width=cell_w - 1.0, height=0.55,
            text=body, size=11, color=INK_SOFT,
        )

    add_textbox(
        slide,
        left=0.6,
        top=6.4,
        width=12.0,
        height=0.4,
        text="Every PR held to this bar — not just the first one.",
        size=14,
        bold=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_feature_001(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Feature 001",
        "policy-loop vertical slice · shipped (990b437 + a49939e)",
        footer="docs/runbook/feature-001-readiness-review.md · every NON-NEGOTIABLE PASS with a named artifact",
    )

    # What it ships
    add_textbox(
        slide, left=0.6, top=1.8, width=12.0, height=0.4,
        text="Ships", size=14, bold=True, color=ACCENT_DARK,
    )
    ships = (
        "4-node LangGraph (Orchestrator · Generator · Validator · Deployer)  ·  "
        "real SLM in CI contract + integration (temp=0, seed=0xC0FFEE)  ·  "
        "RLS-scoped audit_events with FR-017a minimum field set  ·  "
        "immutable policy registry  ·  Collector AI simulator with VSS v6.0 enforcement  ·  "
        "Grafana dashboard + Alertmanager + runbooks  ·  Locust smoke load on PR-tier"
    )
    add_textbox(
        slide, left=0.6, top=2.15, width=12.0, height=1.4, text=ships, size=13, color=INK,
    )

    # ADRs
    add_textbox(
        slide, left=0.6, top=3.65, width=12.0, height=0.4,
        text="Anchoring ADRs", size=14, bold=True, color=ACCENT_DARK,
    )
    adrs = [
        ("ADR-0001", "Pin COVESA VSS to v6.0 (canonical signal vocabulary)."),
        ("ADR-0002", "Default SLM Qwen2.5-7B-Instruct at HF revision SHA a09a35458…  (Proposed; eval baseline gated to GPU runner)."),
        ("ADR-0003", "outlines as the constrained-decoding library."),
        ("ADR-0004", "Deterministic-fingerprint stub for load + soak (Principle XIV split)."),
        ("ADR-0005", "ECS-on-EC2 with g5/g6 GPU node group on AWS."),
        ("ADR-0006", "DevDefaultPolicyClient for local foundation smoke (startup-guarded)."),
    ]
    y = 4.05
    for idx, (a, body) in enumerate(adrs):
        add_textbox(
            slide, left=0.6, top=y, width=1.4, height=0.3,
            text=a, size=11, bold=True, color=ACCENT_DARK, mono=True,
        )
        add_textbox(
            slide, left=2.0, top=y, width=10.5, height=0.3, text=body, size=11, color=INK,
        )
        y += 0.32
        if idx == 2:
            y += 0.02

    return slide


def slide_feature_002(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Feature 002",
        "multi-tenant isolation · shipped (6b46c78)",
        footer="docs/runbook/feature-002-readiness-review.md · all 14 phases closed · ADR-0007 + ADR-0009 promoted Accepted",
    )

    add_textbox(
        slide, left=0.6, top=1.8, width=12.0, height=0.4,
        text="Adds", size=14, bold=True, color=ACCENT_DARK,
    )
    adds = (
        "RESTRICTIVE RLS on every tenant-scoped table + collectmind_tenant role + SET LOCAL ROLE pattern  ·  "
        "break-glass router (distinct audience, atomic kind=break_glass audit row)  ·  "
        "per-tenant token-bucket rate limit (Redis Lua, 3-branch failure-CLOSED)  ·  "
        "hot-store key migration to tenant_id:vehicle_id:signal_name (env-gated dual-read + Fatal deadline)  ·  "
        "deployer-node tenant-scope check (Fatal supersedes Recoverable retry by topology)  ·  "
        "tenant_vehicles + tenant_vehicles_history with immutability trigger"
    )
    add_textbox(
        slide, left=0.6, top=2.15, width=12.0, height=1.7, text=adds, size=13, color=INK,
    )

    add_textbox(
        slide, left=0.6, top=4.0, width=12.0, height=0.4,
        text="Anchoring ADRs", size=14, bold=True, color=ACCENT_DARK,
    )
    adrs = [
        ("ADR-0007", "RESTRICTIVE RLS + PERMISSIVE baseline + break-glass with atomic audit  (Accepted)."),
        ("ADR-0008", "Per-tenant rate limiting + hot-store key migration mechanism  (Proposed; promotes on first workflow_dispatch SC-002+SC-003)."),
        ("ADR-0009", "Tenant-vehicle ownership store: mutable current row + append-only history + write-through Redis cache  (Accepted)."),
    ]
    y = 4.4
    for a, body in adrs:
        add_textbox(
            slide, left=0.6, top=y, width=1.4, height=0.3,
            text=a, size=11, bold=True, color=ACCENT_DARK, mono=True,
        )
        add_textbox(
            slide, left=2.0, top=y, width=10.5, height=0.6,
            text=body, size=11, color=INK,
        )
        y += 0.55

    return slide


def slide_slo(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Engineering rigor — measured, not aspired",
        "Principle XI binds; Principle XIV gates",
        footer="constitution Principles XI + XIV · docs/PROJECT_STATE.md (recorded measurements section)",
    )

    add_textbox(
        slide, left=0.6, top=1.8, width=12.0, height=0.4,
        text="Measured (local runs · no fabrication)", size=14, bold=True, color=ACCENT_DARK,
    )

    measured = [
        ("SC-006 dashboard lag", "2.11 s", "10 s", "~5×", "T136 · max of 5 runs"),
        ("SC-008 quickstart e2e", "3 s", "600 s", "~200×", "T292 · warm Compose"),
        ("Coverage (Principle IV)", "85.36 %", "≥ 85 %", "+0.36 pp", "T285 · across-tier"),
        ("Smoke load p50", "50 ms", "4 s", "~80×", "T134 · 280 req · 0 fail"),
        ("CI wall-clock (SC-009)", "11m 02s", "20 min", "9 min", "PR #1 · run 25775623611"),
    ]
    y = 2.2
    for label, val, budget, headroom, src in measured:
        add_panel(slide, left=0.6, top=y, width=12.0, height=0.42)
        add_textbox(slide, left=0.75, top=y + 0.07, width=3.6, height=0.32,
                    text=label, size=11, bold=True, color=INK)
        add_textbox(slide, left=4.4, top=y + 0.07, width=1.6, height=0.32,
                    text=val, size=12, bold=True, color=OK, mono=True)
        add_textbox(slide, left=6.05, top=y + 0.07, width=1.4, height=0.32,
                    text="budget " + budget, size=10, color=INK_SOFT)
        add_textbox(slide, left=7.5, top=y + 0.07, width=1.3, height=0.32,
                    text=headroom, size=10, bold=True, color=OK)
        add_textbox(slide, left=8.85, top=y + 0.07, width=3.6, height=0.32,
                    text=src, size=9, color=INK_MUTED)
        y += 0.48

    add_textbox(
        slide, left=0.6, top=5.0, width=12.0, height=0.4,
        text="Gated (workflow_dispatch · nightly · GPU runner)", size=14, bold=True, color=WARN,
    )
    gated = [
        ("SC-002", "1000 ev/s/tenant · 30 min · ≥ 99.9 %", "workflow_dispatch · locust quitting hook"),
        ("SC-003", "24h soak · mem growth ≤ 5 % · err ≤ 0.1 %", "nightly · post-run RSS gate"),
        ("ADR-0002 eval baseline", "VSS pass rate + p50 gen latency", "[self-hosted, gpu] · eval-suite job · bracketed"),
    ]
    y = 5.4
    for label, body, src in gated:
        add_panel(slide, left=0.6, top=y, width=12.0, height=0.42)
        add_textbox(slide, left=0.75, top=y + 0.07, width=3.8, height=0.32,
                    text=label, size=11, bold=True, color=WARN, mono=True)
        add_textbox(slide, left=4.6, top=y + 0.07, width=4.4, height=0.32,
                    text=body, size=11, color=INK)
        add_textbox(slide, left=9.05, top=y + 0.07, width=3.4, height=0.32,
                    text=src, size=9, color=INK_MUTED)
        y += 0.48

    add_textbox(
        slide, left=0.6, top=6.85, width=12.0, height=0.3,
        text="What's not measured here is gated, not skipped. The gating condition is named.",
        size=10, color=INK_MUTED, mono=True,
    )

    return slide


def slide_audit(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Audit is structural",
        "FR-017a minimum field set on every kind=generated row",
        footer="constitution Principle XVII · src/collectmind/registry/audit.py · raises ValueError if any FR-017a field missing",
    )

    add_textbox(
        slide, left=0.6, top=1.75, width=12.0, height=0.4,
        text="One row from the demo's audit chain (recorded mode, tenant-a)",
        size=12, color=INK_SOFT,
    )

    audit_lines = [
        ('"event_id":              "01HW1A2B3C4D5E6F7G8H9J0KQA-02"',),
        ('"kind":                  "generated"',),
        ('"occurred_at":           "2026-05-13T12:00:01.347Z"',),
        ('"correlation_id":        "01HW1A2B3C4D5E6F7G8H9J0KQA"',),
        ('"principal_subject":     "graph/policy-generator"',),
        ('"slm_repo":              "Qwen/Qwen2.5-7B-Instruct"',),
        ('"slm_revision_sha":      "a09a35458c702b33eeacc393d103063234e8bc28"',),
        ('"slm_runtime":           "vllm"',),
        ('"slm_runtime_version":   "0.20.1"',),
        ('"slm_quantization":      "bf16"',),
        ('"slm_decoding_seed":     12648430',),
        ('"prompt_template_version": "1.0.0"',),
        ('"policy_ref": {',),
        ('   "tenant_id": "tenant-a",',),
        ('   "policy_id": "pol-tenant-a-brake-wear",',),
        ('   "version":   "1.0.0"',),
        ('}',),
    ]

    add_panel(slide, left=0.6, top=2.2, width=12.0, height=4.0,
              fill=RGBColor(0xFA, 0xFA, 0xFA), border=PANEL_BORDER)
    y = 2.32
    for (line,) in audit_lines:
        add_textbox(
            slide, left=0.85, top=y, width=11.5, height=0.22,
            text=line, size=10, color=INK, mono=True,
        )
        y += 0.22

    add_textbox(
        slide, left=0.6, top=6.4, width=12.0, height=0.4,
        text="Audit is not logging. Audit is the contract — queryable by API, immutable at schema level.",
        size=13, bold=True, color=INK, align=PP_ALIGN.CENTER,
    )

    return slide


def slide_multitenant(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Multi-tenant isolation",
        "RESTRICTIVE RLS by default · break-glass is a distinct router",
        footer="ADR-0007 · src/collectmind/registry/db.py (Database.acquire) · audit-admin.v1.yaml",
    )

    # Two-column compare
    col_w = 5.8
    col_top = 1.8
    col_h = 4.6
    add_panel(slide, left=0.6, top=col_top, width=col_w, height=col_h)
    add_panel(slide, left=6.95, top=col_top, width=col_w, height=col_h, fill=RGBColor(0xEE, 0xE9, 0xFE))

    # Left column: tenant path
    add_textbox(slide, left=0.85, top=col_top + 0.15, width=col_w - 0.4, height=0.4,
                text="Tenant path (default)", size=14, bold=True, color=INK)
    items_l = [
        ("audience", "collectmind-tenant"),
        ("DB role", "collectmind_tenant (non-BYPASSRLS)"),
        ("RLS policy", "RESTRICTIVE per-tenant + PERMISSIVE baseline"),
        ("Cross-tenant access", "HTTP 404 (FR-006: no existence oracle)"),
        ("Code path", "src/collectmind/auth/dependencies.py → Database.acquire(tenant_id)"),
        ("Verification", "tests/integration/test_rls_restrictive.py · 3-assertion contract"),
        ("Cost on bypass", "structurally impossible — role is non-BYPASSRLS"),
    ]
    y = col_top + 0.65
    for k, v in items_l:
        add_textbox(slide, left=0.85, top=y, width=col_w - 0.4, height=0.22,
                    text=k.upper(), size=8, bold=True, color=ACCENT_DARK)
        add_textbox(slide, left=0.85, top=y + 0.22, width=col_w - 0.4, height=0.38,
                    text=v, size=11, color=INK)
        y += 0.6

    # Right column: break-glass path
    add_textbox(slide, left=7.20, top=col_top + 0.15, width=col_w - 0.4, height=0.4,
                text="Break-glass path (operator-only)", size=14, bold=True, color=INK)
    items_r = [
        ("audience", "collectmind-operator"),
        ("DB primitive", "service-principal connection (BYPASSRLS)"),
        ("RLS policy", "bypassed — parameterized on tenant_scope, not widenable mid-flight"),
        ("Atomic audit", "kind=break_glass row inside same txn as bypassed SELECT (FR-005b)"),
        ("Code path", "src/collectmind/audit_admin/api.py — distinct router, shares no code with the tenant path"),
        ("Reason codes", "5 enumerated · extension requires ADR-0007 amendment"),
        ("Alerting", "BreakGlassInvoked page-tier + BreakGlassBurstInvocation critical-tier (per operator)"),
    ]
    y = col_top + 0.65
    for k, v in items_r:
        add_textbox(slide, left=7.20, top=y, width=col_w - 0.4, height=0.22,
                    text=k.upper(), size=8, bold=True, color=ACCENT_DARK)
        add_textbox(slide, left=7.20, top=y + 0.22, width=col_w - 0.4, height=0.38,
                    text=v, size=11, color=INK)
        y += 0.6

    return slide


def slide_cadence(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "The cadence",
        "Reproducible methodology: two features shipped through this loop",
        footer="Spec-Kit · constitution.md is the highest-priority artifact · overrides every plan choice in conflict",
    )

    steps = [
        ("1", "constitution", "highest priority"),
        ("2", "ADR(s)", "non-obvious decisions"),
        ("3", "/speckit-specify", "what + why"),
        ("4", "/speckit-clarify", "underspecified bits"),
        ("5", "/speckit-plan", "tech + structure"),
        ("6", "/speckit-tasks", "dependency-ordered"),
        ("7", "red commits", "tests first"),
        ("8", "/speckit-implement", "drive green"),
        ("9", "readiness review", "every NN PASS"),
    ]

    n = len(steps)
    box_w = 1.25
    gap = 0.10
    start_x = 0.6
    top = 2.6
    for i, (num, label, sub) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_panel(slide, left=x, top=top, width=box_w, height=1.5)
        add_textbox(
            slide, left=x, top=top + 0.15, width=box_w, height=0.35,
            text=num, size=18, bold=True, color=ACCENT_DARK, mono=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left=x + 0.08, top=top + 0.55, width=box_w - 0.16, height=0.4,
            text=label, size=10, bold=True, color=INK, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left=x + 0.08, top=top + 0.97, width=box_w - 0.16, height=0.5,
            text=sub, size=8, color=INK_SOFT, align=PP_ALIGN.CENTER,
        )
        # arrow to next
        if i < n - 1:
            ax = x + box_w + 0.005
            add_textbox(
                slide, left=ax, top=top + 0.55, width=gap, height=0.5,
                text="→", size=14, bold=True, color=INK_FAINT, align=PP_ALIGN.CENTER,
            )

    add_textbox(
        slide, left=0.6, top=4.6, width=12.0, height=0.4,
        text="Feature 001 closed at 990b437 + a49939e · feature 002 closed at 6b46c78 · same loop both times.",
        size=12, color=INK, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, left=0.6, top=5.1, width=12.0, height=0.4,
        text="Feature 003 will run the same cadence. Methodology is reproducible, not improvised.",
        size=12, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER,
    )

    return slide


def slide_deferred(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "What's deferred and why",
        "Deferrals are named, not silent · each gated to a real trigger",
        footer="docs/PROJECT_STATE.md → \"What is next — Phase 7\" table · constitution Principle XI no-fabrication rule",
    )

    rows = [
        (
            "ADR-0002 eval baseline + Accepted promotion",
            "GPU runner not yet attached to the repo",
            "First successful workflow_dispatch eval-suite on [self-hosted, gpu]",
        ),
        (
            "ADR-0008 promotion to Accepted",
            "Local-stack covers structural primitives; production-rate verification is the missing piece",
            "First workflow_dispatch SC-002 + SC-003 runs against the rate-limited orchestration-api",
        ),
        (
            "SC-009 rolling-5-PR aggregator",
            "Inaugural run at 11m 02s — aggregator not needed at the inaugural measurement",
            "First PR-tier run trending toward the 18-min warning threshold",
        ),
        (
            "Supply-chain refresh sweep",
            "Trivy + pip-audit + Syft cataloger config + base image + 8 pinned-dep bumps",
            "Deliberate dependency-update sprint; PR titled chore: supply-chain refresh",
        ),
        (
            "test_rls_migration_rollback schema_migrations desync",
            "Pre-existing test-infra flake; rollback helpers skip the tracking table",
            "Two-line fix; clear schema_migrations rows in the rollback helper",
        ),
        (
            "Unit-tier coverage parity for feature-002 routers",
            "Feature-002 routers reachable only through contract/integration tiers; unit cov ≈ 83 %",
            "Add mocked unit-tier tests (TestClient + AsyncMock); re-tighten unit gate to 85",
        ),
    ]

    y = 1.85
    for label, reason, gate in rows:
        add_panel(slide, left=0.6, top=y, width=12.0, height=0.75)
        add_textbox(slide, left=0.78, top=y + 0.08, width=4.2, height=0.6,
                    text=label, size=11, bold=True, color=INK)
        add_textbox(slide, left=5.05, top=y + 0.08, width=3.6, height=0.6,
                    text=reason, size=10, color=INK_SOFT)
        add_textbox(slide, left=8.75, top=y + 0.08, width=3.8, height=0.6,
                    text="GATE: " + gate, size=10, color=ACCENT_DARK, mono=True)
        y += 0.82

    return slide


def slide_demo_url(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Try it",
        "Deployed UI · recorded mode · public URL",
        footer="GitHub Pages · same artifact builds for Vercel (Q2 build-prompt decision) · cd demo && vercel --prod",
    )

    add_panel(slide, left=0.6, top=1.9, width=12.0, height=2.0, fill=PANEL)
    add_textbox(
        slide, left=0.85, top=2.05, width=11.5, height=0.4,
        text="DEMO", size=10, bold=True, color=ACCENT_DARK,
    )
    add_textbox(
        slide, left=0.85, top=2.4, width=11.5, height=0.6,
        text="arunveligatla99.github.io/snts_collectmind",
        size=24, bold=True, color=INK, mono=True,
    )
    add_textbox(
        slide, left=0.85, top=3.05, width=11.5, height=0.7,
        text="Recorded-mode build (deployed-only posture; no public Compose, no shipped tokens). "
             "Six routes: Overview · Operator · Audit · SLO · Tenants · Docs.",
        size=12, color=INK_SOFT,
    )

    add_panel(slide, left=0.6, top=4.15, width=12.0, height=2.0, fill=PANEL)
    add_textbox(
        slide, left=0.85, top=4.30, width=11.5, height=0.4,
        text="REPO", size=10, bold=True, color=ACCENT_DARK,
    )
    add_textbox(
        slide, left=0.85, top=4.65, width=11.5, height=0.6,
        text="github.com/arunveligatla99/snts_collectmind",
        size=24, bold=True, color=INK, mono=True,
    )
    add_textbox(
        slide, left=0.85, top=5.30, width=11.5, height=0.7,
        text="Branch arun/develop carries feature 001 + 002 shipped. "
             "Branch arun/demo-ui carries the demo UI. Constitution at .specify/memory/constitution.md.",
        size=12, color=INK_SOFT,
    )

    add_textbox(
        slide, left=0.6, top=6.40, width=12.0, height=0.4,
        text="make demo  ·  make demo-test  ·  make demo-build  ·  bash demo/scripts/record_fixtures.sh",
        size=11, color=INK_MUTED, mono=True, align=PP_ALIGN.CENTER,
    )

    return slide


def slide_why_sonatus(prs):
    slide = new_slide(prs)
    add_slide_title(
        slide,
        "Why I'm interested in Sonatus",
        "[fill in your own words]",
        footer="placeholder · personal narrative · one slide · no boilerplate",
    )

    add_panel(slide, left=0.6, top=2.0, width=12.0, height=4.0, fill=PANEL)
    add_textbox(
        slide, left=0.85, top=2.4, width=11.5, height=0.6,
        text="[Your narrative goes here.]",
        size=20, bold=True, color=INK_FAINT, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, left=0.85, top=3.3, width=11.5, height=2.5,
        text=(
            "Anchor points if useful: AI Director's edge-AI + SLM thesis · OEM-customer "
            "audit and compliance posture · the in-vehicle vs. cloud split · production-grade "
            "vs. demo-grade engineering rigor.\n\n"
            "Replace this text with one paragraph in your own voice."
        ),
        size=13, color=INK_SOFT, align=PP_ALIGN.CENTER,
    )

    return slide


def slide_qa(prs):
    slide = new_slide(prs)
    add_accent_bar(slide)
    add_textbox(
        slide, left=0.6, top=1.4, width=12.0, height=0.8,
        text="Q & A",
        size=64, bold=True, color=INK,
    )
    add_textbox(
        slide, left=0.6, top=2.6, width=12.0, height=0.5,
        text="Happy to deep-dive any layer.",
        size=20, color=INK_SOFT,
    )

    add_panel(slide, left=0.6, top=4.2, width=12.0, height=2.4)
    add_textbox(slide, left=0.85, top=4.4, width=11.5, height=0.4,
                text="Arun Veligatla", size=18, bold=True, color=INK)
    add_textbox(slide, left=0.85, top=4.85, width=11.5, height=0.4,
                text="arun.veligatla@gmail.com", size=14, color=INK, mono=True)
    add_textbox(slide, left=0.85, top=5.25, width=11.5, height=0.4,
                text="github.com/arunveligatla99", size=14, color=INK, mono=True)
    add_textbox(slide, left=0.85, top=5.70, width=11.5, height=0.4,
                text="arunveligatla99.github.io/snts_collectmind", size=14, color=ACCENT_DARK, mono=True)

    return slide


def build():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_gap,
        slide_what_it_does,
        slide_architecture,
        slide_constitution,
        slide_feature_001,
        slide_feature_002,
        slide_slo,
        slide_audit,
        slide_multitenant,
        slide_cadence,
        slide_deferred,
        slide_demo_url,
        slide_why_sonatus,
        slide_qa,
    ]
    for b in builders:
        b(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT.relative_to(REPO_ROOT)} ({len(builders)} slides)")


if __name__ == "__main__":
    build()
